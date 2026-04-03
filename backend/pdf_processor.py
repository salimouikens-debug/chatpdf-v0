import uuid
from io import BytesIO

from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from pinecone import Pinecone

from config import (
    OPENAI_API_KEY,
    PINECONE_API_KEY,
    PINECONE_INDEX_NAME,
    CHUNK_SIZE,
    CHUNK_OVERLAP,
    EMBEDDING_MODEL,
)
from database import db_save_document, db_get_document, db_list_documents, db_delete_document

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    length_function=len,
    separators=["\n\n", "\n", ". ", " ", ""],
)

_embeddings = None
_index = None


def _get_embeddings():
    global _embeddings
    if _embeddings is None:
        if not OPENAI_API_KEY:
            raise RuntimeError("OPENAI_API_KEY is not set in .env.local")
        _embeddings = OpenAIEmbeddings(model=EMBEDDING_MODEL, openai_api_key=OPENAI_API_KEY)
    return _embeddings


def _get_index():
    global _index
    if _index is None:
        if not PINECONE_API_KEY:
            raise RuntimeError("PINECONE_API_KEY is not set in .env.local")
        pc = Pinecone(api_key=PINECONE_API_KEY)
        _index = pc.Index(PINECONE_INDEX_NAME)
    return _index


import os

def extract_text_by_page(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract text from each page/slide/section, returning a list of {page, text}."""
    ext = os.path.splitext(filename)[1].lower()
    pages = []
    
    if ext == ".pdf":
        reader = PdfReader(BytesIO(file_bytes))
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append({"page": i + 1, "text": text})
    elif ext in [".docx", ".doc"]:
        doc = docx.Document(BytesIO(file_bytes))
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if text.strip():
            pages.append({"page": 1, "text": text.strip()})
    elif ext == ".pptx":
        prs = Presentation(BytesIO(file_bytes))
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            text = "\n".join(slide_text).strip()
            if text:
                pages.append({"page": i + 1, "text": text})
    elif ext == ".txt":
        text = file_bytes.decode('utf-8', errors='ignore').strip()
        if text:
            pages.append({"page": 1, "text": text})

    return pages


def chunk_pages(pages: list[dict]) -> list[dict]:
    """Split page texts into smaller chunks, preserving page metadata."""
    chunks = []
    for page_info in pages:
        page_chunks = text_splitter.split_text(page_info["text"])
        for chunk_text in page_chunks:
            chunks.append({"page": page_info["page"], "text": chunk_text})
    return chunks


async def process_pdf(file_bytes: bytes, filename: str, user_id: str) -> dict:
    """Full pipeline: extract -> chunk -> embed -> store in Pinecone + PostgreSQL."""
    doc_id = str(uuid.uuid4())

    pages = extract_text_by_page(file_bytes, filename)
    if not pages:
        raise ValueError(f"Could not extract any text from the document {filename}.")

    total_pages = max(p["page"] for p in pages)
    chunks = chunk_pages(pages)

    embeddings = _get_embeddings()
    texts = [c["text"] for c in chunks]
    vectors = await embeddings.aembed_documents(texts)

    index = _get_index()
    upsert_data = []
    for i, (chunk, vector) in enumerate(zip(chunks, vectors)):
        vec_id = f"{doc_id}#{i}"
        metadata = {
            "document_id": doc_id,
            "filename": filename,
            "page": chunk["page"],
            "text": chunk["text"],
        }
        upsert_data.append((vec_id, vector, metadata))

    batch_size = 100
    for i in range(0, len(upsert_data), batch_size):
        batch = upsert_data[i : i + batch_size]
        index.upsert(vectors=batch)

    db_save_document(user_id, doc_id, filename, total_pages, len(chunks), pdf_data=file_bytes)

    return {
        "id": doc_id,
        "filename": filename,
        "total_pages": total_pages,
        "total_chunks": len(chunks),
    }


def delete_document(doc_id: str) -> bool:
    """Delete all vectors for a document from Pinecone + SQLite."""
    doc = db_get_document(doc_id)
    if not doc:
        return False

    index = _get_index()
    vec_ids = [f"{doc_id}#{i}" for i in range(doc["total_chunks"])]
    batch_size = 100
    for i in range(0, len(vec_ids), batch_size):
        batch = vec_ids[i : i + batch_size]
        index.delete(ids=batch)

    db_delete_document(doc_id)
    return True


def list_documents(user_id: str) -> list[dict]:
    """Return all uploaded document metadata for a given user."""
    return db_list_documents(user_id)
